"""
PPTX Release Notes generator (SE-branded, PICASso format).

Pipeline:
  1. Load content from JSON spec file
  2. Build PPTX presentation using SE-branded slide templates
  3. Write output to output/full_release_notes/

Slide structure (matching PICASso template):
  1. Title Slide - product name, version, release date
  2. Executive Summary - key highlights
  3-N. Feature Detail - one slide per feature section
  N+1. Defect Fixes - checkmark list
  N+2. Key Benefits - summary benefit cards
"""

from __future__ import annotations

import json
from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from src.config.settings import Settings
from src.utils.date_utils import today_str
from src.utils.file_utils import ensure_dir, sanitize_filename

# ---------------------------------------------------------------------------
# Color palette (Schneider Electric brand)
# ---------------------------------------------------------------------------

SE_GREEN = RGBColor(0x3D, 0xCD, 0x58)
SE_DARK_GREEN = RGBColor(0x00, 0x95, 0x2C)
SE_DARKER_GREEN = RGBColor(0x00, 0x6B, 0x1E)
SE_LIGHT_GREEN_BG = RGBColor(0xED, 0xF9, 0xF0)
SE_LIGHT_GREEN_BG2 = RGBColor(0xEB, 0xF7, 0xED)
SE_SEPARATOR = RGBColor(0xD6, 0xEC, 0xDB)
SE_TEXT_DARK = RGBColor(0x33, 0x33, 0x33)
SE_TEXT_GRAY = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT_NAME = "Calibri"

# Slide dimensions (widescreen 16:9)
SLIDE_WIDTH = Emu(12192000)   # 13.333 inches
SLIDE_HEIGHT = Emu(6858000)   # 7.500 inches


# ---------------------------------------------------------------------------
# Data models
# ---------------------------------------------------------------------------

@dataclass
class BulletItem:
    """A bullet point with optional secondary detail text."""
    text: str
    detail: str = ""


@dataclass
class FeatureSlide:
    """Content for a single feature detail slide."""
    title: str
    subtitle: str
    bullets: list[BulletItem] = field(default_factory=list)
    screenshot_caption: str = ""


@dataclass
class ExecutiveSummaryItem:
    """One numbered highlight on the executive summary slide."""
    title: str
    description: str


@dataclass
class KeyBenefit:
    """One benefit card on the key benefits slide."""
    number: str       # "01", "02", etc.
    category: str     # "COMPLIANCE & AUDIT"
    title: str
    description: str


@dataclass
class DefectFixesData:
    """Content for the defect fixes slide."""
    subtitle: str
    fixes: list[str] = field(default_factory=list)


@dataclass
class PptxContent:
    """Complete content model for a PPTX release notes presentation."""
    product_name: str
    version_display: str   # "Release 9.0"
    version_code: str      # "PIC-2026-RC-9.0"
    release_date: str      # "05 March 2026"
    executive_summary: list[ExecutiveSummaryItem] = field(default_factory=list)
    features: list[FeatureSlide] = field(default_factory=list)
    defect_fixes: DefectFixesData = field(default_factory=lambda: DefectFixesData(subtitle=""))
    key_benefits: list[KeyBenefit] = field(default_factory=list)


# ---------------------------------------------------------------------------
# Slide builder (internal)
# ---------------------------------------------------------------------------

class _SlideBuilder:
    """Builds SE-branded PPTX slides matching the PICASso release notes template."""

    def __init__(self) -> None:
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_WIDTH
        self.prs.slide_height = SLIDE_HEIGHT
        self._blank = self.prs.slide_layouts[6]  # Blank layout

    # -- low-level helpers --------------------------------------------------

    def _slide(self):
        return self.prs.slides.add_slide(self._blank)

    def _rect(self, slide, left, top, width, height, color):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height,
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    def _oval(self, slide, left, top, size, color):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, left, top, size, size,
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    def _text(self, slide, text, left, top, width, height,
              size=Pt(12), color=SE_TEXT_DARK, bold=False, italic=False,
              align=PP_ALIGN.LEFT, wrap=True, anchor=MSO_ANCHOR.TOP):
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = wrap
        tf.auto_size = None
        try:
            tf.vertical_anchor = anchor
        except Exception:
            pass
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.name = FONT_NAME
        run.font.size = size
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        return box

    # -- shared slide elements ----------------------------------------------

    def _footer(self, slide):
        """Green footer bar at bottom of every slide."""
        self._rect(slide, Emu(0), Inches(7.36), SLIDE_WIDTH, Inches(0.14), SE_GREEN)
        self._text(slide, "Confidential \u2014 For Internal Use",
                   Inches(0.42), Inches(7.36), Inches(6.0), Inches(0.14),
                   size=Pt(7), color=WHITE)
        self._text(slide, "Schneider Electric",
                   Inches(10.33), Inches(7.36), Inches(2.80), Inches(0.14),
                   size=Pt(7), color=WHITE, bold=True, align=PP_ALIGN.RIGHT)

    def _content_header(self, slide, title: str, version_code: str):
        """Green header bar with title and version badge (slides 2+)."""
        # Header bar
        self._rect(slide, Emu(0), Emu(0), SLIDE_WIDTH, Inches(0.80), SE_GREEN)
        # Dark green accent strip
        self._rect(slide, Emu(0), Emu(0), Inches(0.20), Inches(0.80), SE_DARK_GREEN)
        # Title
        self._text(slide, title,
                   Inches(0.34), Inches(0.13), Inches(10.13), Inches(0.60),
                   size=Pt(24), color=WHITE, bold=True)
        # Version badge background
        self._rect(slide, Inches(10.95), Inches(0.23), Inches(2.10), Inches(0.34), SE_DARK_GREEN)
        # Version badge text
        self._text(slide, version_code,
                   Inches(10.95), Inches(0.23), Inches(2.10), Inches(0.34),
                   size=Pt(10), color=WHITE, bold=True, align=PP_ALIGN.CENTER, wrap=False)

    def _subtitle_band(self, slide, text: str):
        """Light green subtitle band below header."""
        self._rect(slide, Emu(0), Inches(0.80), SLIDE_WIDTH, Inches(0.46), SE_LIGHT_GREEN_BG)
        self._rect(slide, Emu(0), Inches(0.80), Inches(0.07), Inches(0.46), SE_GREEN)
        self._text(slide, text,
                   Inches(0.18), Inches(0.87), Inches(12.97), Inches(0.34),
                   size=Pt(12), color=SE_DARKER_GREEN, italic=True, wrap=False)

    # -- slide builders -----------------------------------------------------

    def build_title_slide(self, content: PptxContent):
        """Slide 1: Cover page with green left panel."""
        slide = self._slide()

        # Green left panel
        self._rect(slide, Emu(0), Emu(0), Inches(5.60), SLIDE_HEIGHT, SE_GREEN)
        # Dark green accent strip
        self._rect(slide, Emu(0), Emu(0), Inches(0.20), SLIDE_HEIGHT, SE_DARK_GREEN)
        # Dark accent bar at top
        self._rect(slide, Inches(0.32), Inches(0.34), Inches(4.00), Inches(0.30), SE_DARKER_GREEN)

        # "RELEASE NOTES" label
        self._text(slide, "RELEASE NOTES",
                   Inches(0.36), Inches(0.90), Inches(5.10), Inches(0.36),
                   size=Pt(11), color=WHITE, bold=True)

        # Product name + version (large)
        self._text(slide, f"{content.product_name}\n{content.version_display}",
                   Inches(0.36), Inches(1.35), Inches(5.05), Inches(2.00),
                   size=Pt(38), color=WHITE, bold=True)

        # White divider line
        self._rect(slide, Inches(0.36), Inches(3.55), Inches(4.60), Inches(0.02), WHITE)

        # Version code + release date
        self._text(slide, f"{content.version_code}  \u00b7  Move to Prod: {content.release_date}",
                   Inches(0.36), Inches(3.75), Inches(5.05), Inches(0.50),
                   size=Pt(14), color=WHITE)

        # "Schneider Electric" on white side
        self._text(slide, "Schneider Electric",
                   Inches(5.90), Inches(6.85), Inches(7.00), Inches(0.46),
                   size=Pt(15), color=SE_DARK_GREEN, bold=True)

        # Footer
        self._footer(slide)

    def build_executive_summary(self, content: PptxContent):
        """Slide 2: Numbered key highlights."""
        slide = self._slide()
        self._content_header(slide, f"{content.version_display} \u2014 Executive Summary", content.version_code)
        self._footer(slide)

        items = content.executive_summary
        if not items:
            return

        # Vertical layout: each item gets ~1.5in, starting at 1.12in
        y_start = 1.12
        item_height = 1.54  # distance between items

        for i, item in enumerate(items):
            y = y_start + i * item_height

            # Green number circle
            self._oval(slide, Inches(0.42), Inches(y), Inches(0.48), SE_GREEN)
            self._text(slide, str(i + 1),
                       Inches(0.50), Inches(y + 0.07), Inches(0.32), Inches(0.37),
                       size=Pt(22), color=WHITE, bold=True, align=PP_ALIGN.CENTER)

            # Title
            self._text(slide, item.title,
                       Inches(1.04), Inches(y), Inches(6.69), Inches(0.40),
                       size=Pt(18), color=SE_TEXT_DARK, bold=True)

            # Description
            self._text(slide, item.description,
                       Inches(1.04), Inches(y + 0.44), Inches(7.67), Inches(0.34),
                       size=Pt(12), color=SE_TEXT_GRAY)

            # Separator line (except after last item)
            if i < len(items) - 1:
                sep_y = y + item_height - 0.11
                self._rect(slide, Inches(0.42), Inches(sep_y), Inches(7.31), Inches(0.01), SE_SEPARATOR)

    def build_feature_slide(self, feature: FeatureSlide, version_code: str):
        """Slides 3-N: Feature detail with bullets and screenshot placeholder."""
        slide = self._slide()
        self._content_header(slide, feature.title, version_code)
        self._subtitle_band(slide, feature.subtitle)
        self._footer(slide)

        bullets = feature.bullets
        if not bullets:
            return

        # Determine layout: bullets with details get more vertical space
        has_details = any(b.detail for b in bullets)
        content_left = 0.66
        content_width = 5.56
        y = 1.44

        if has_details:
            # Paired layout: primary (0.54in) + secondary (0.30in) + gap
            group_height = 1.06
            for b in bullets:
                if y > 6.4:
                    break
                # Green bullet marker
                self._oval(slide, Inches(0.44), Inches(y + 0.16), Inches(0.13), SE_GREEN)
                # Primary text
                self._text(slide, b.text,
                           Inches(content_left), Inches(y), Inches(content_width), Inches(0.54),
                           size=Pt(14), color=SE_TEXT_DARK)
                # Secondary text
                if b.detail:
                    self._text(slide, b.detail,
                               Inches(content_left), Inches(y + 0.54), Inches(content_width), Inches(0.30),
                               size=Pt(10), color=SE_TEXT_GRAY)
                y += group_height
        else:
            # Single-line layout: tighter spacing
            line_height = 0.54
            sub_height = 0.30
            for b in bullets:
                if y > 6.4:
                    break
                self._oval(slide, Inches(0.44), Inches(y + 0.16), Inches(0.13), SE_GREEN)
                self._text(slide, b.text,
                           Inches(content_left), Inches(y), Inches(content_width), Inches(line_height),
                           size=Pt(14), color=SE_TEXT_DARK)
                y += line_height + sub_height

        # Screenshot placeholder on the right
        scr_left = 6.46
        scr_top = 1.44
        scr_width = 6.45
        scr_height = 5.82

        # Light green background rectangle
        self._rect(slide, Inches(scr_left), Inches(scr_top),
                   Inches(scr_width), Inches(scr_height), SE_LIGHT_GREEN_BG2)

        # Crosshair in center
        cx = scr_left + scr_width / 2
        cy = scr_top + scr_height * 0.45
        self._rect(slide, Inches(cx - 0.32), Inches(cy), Inches(0.64), Inches(0.03), SE_DARK_GREEN)
        self._rect(slide, Inches(cx), Inches(cy - 0.32), Inches(0.03), Inches(0.64), SE_DARK_GREEN)

        # "[ Screenshot ]" label
        self._text(slide, "[ Screenshot ]",
                   Inches(scr_left + 0.10), Inches(cy + 0.50), Inches(scr_width - 0.20), Inches(0.34),
                   size=Pt(14), color=SE_TEXT_GRAY, align=PP_ALIGN.CENTER)

        # Caption
        if feature.screenshot_caption:
            self._text(slide, feature.screenshot_caption,
                       Inches(scr_left + 0.15), Inches(cy + 0.82), Inches(scr_width - 0.30), Inches(1.10),
                       size=Pt(10), color=SE_TEXT_GRAY, italic=True)

    def build_defect_fixes_slide(self, defects: DefectFixesData, version_code: str):
        """Defect fixes slide with checkmark list."""
        title = f"{len(defects.fixes)} Defect Fixes \u2014 Quality & Stability"
        slide = self._slide()
        self._content_header(slide, title, version_code)
        self._subtitle_band(slide, defects.subtitle)
        self._footer(slide)

        if not defects.fixes:
            return

        y = 1.54
        row_height = 0.52

        for fix in defects.fixes:
            if y > 6.6:
                break
            # Green square background
            self._rect(slide, Inches(0.42), Inches(y), Inches(0.28), Inches(0.28), SE_GREEN)
            # Checkmark
            self._text(slide, "\u2713",
                       Inches(0.42), Inches(y - 0.02), Inches(0.28), Inches(0.30),
                       size=Pt(14), color=WHITE, bold=True, align=PP_ALIGN.CENTER)
            # Fix description
            self._text(slide, fix,
                       Inches(0.78), Inches(y - 0.03), Inches(12.13), Inches(0.44),
                       size=Pt(13), color=SE_TEXT_DARK)

            y += row_height

    def build_key_benefits_slide(self, content: PptxContent):
        """Key benefits slide with 2x2 card grid + optional bottom bar."""
        slide = self._slide()
        self._content_header(slide, f"Key Benefits \u2014 {content.version_display}", content.version_code)
        self._footer(slide)

        benefits = content.key_benefits
        if not benefits:
            return

        # Card positions: 2x2 grid
        positions = [
            (0.42, 1.02),   # top-left
            (6.78, 1.02),   # top-right
            (0.42, 3.87),   # bottom-left
            (6.78, 3.87),   # bottom-right
        ]
        card_w = 6.13
        card_h = 2.65
        header_h = 1.10

        for i, benefit in enumerate(benefits[:4]):
            if i >= len(positions):
                break
            x, y = positions[i]

            # Card background (white with shadow effect)
            self._rect(slide, Inches(x), Inches(y), Inches(card_w), Inches(card_h), WHITE)
            # Card header (dark green)
            self._rect(slide, Inches(x), Inches(y), Inches(card_w), Inches(header_h), SE_DARK_GREEN)
            # Green accent strip on left of header
            self._rect(slide, Inches(x), Inches(y), Inches(0.12), Inches(header_h), SE_GREEN)

            # Number
            self._text(slide, benefit.number,
                       Inches(x + 0.20), Inches(y + 0.06), Inches(1.00), Inches(1.00),
                       size=Pt(36), color=WHITE, bold=True)

            # Category label (right-aligned in header)
            cat_x = x + card_w - 2.50
            self._text(slide, benefit.category,
                       Inches(cat_x), Inches(y + 0.12), Inches(2.31), Inches(0.37),
                       size=Pt(10), color=WHITE, bold=True, align=PP_ALIGN.RIGHT)

            # Green separator
            self._rect(slide, Inches(x), Inches(y + header_h), Inches(card_w), Inches(0.03), SE_GREEN)

            # Title
            self._text(slide, benefit.title,
                       Inches(x + 0.18), Inches(y + 1.26), Inches(card_w - 0.30), Inches(0.52),
                       size=Pt(14), color=SE_TEXT_DARK, bold=True)

            # Description
            self._text(slide, benefit.description,
                       Inches(x + 0.18), Inches(y + 1.82), Inches(card_w - 0.30), Inches(0.73),
                       size=Pt(10), color=SE_TEXT_GRAY)

        # 5th benefit as bottom bar (if present)
        if len(benefits) >= 5:
            b5 = benefits[4]
            bar_y = 6.68
            bar_h = 0.54
            bar_w = 12.49

            # Dark green bar
            self._rect(slide, Inches(0.42), Inches(bar_y), Inches(bar_w), Inches(bar_h), SE_DARK_GREEN)
            # Green accent on left
            self._rect(slide, Inches(0.42), Inches(bar_y), Inches(0.50), Inches(bar_h), SE_GREEN)

            # Number/icon
            self._text(slide, b5.number,
                       Inches(0.42), Inches(bar_y + 0.09), Inches(0.50), Inches(bar_h),
                       size=Pt(16), color=WHITE, bold=True, align=PP_ALIGN.CENTER)

            # Category
            self._text(slide, b5.category,
                       Inches(0.98), Inches(bar_y + 0.07), Inches(1.30), Inches(0.37),
                       size=Pt(12), color=WHITE, bold=True)

            # Description (separator + text)
            self._rect(slide, Inches(3.22), Inches(bar_y + 0.12), Inches(0.02), Inches(0.30), SE_GREEN)
            self._text(slide, b5.description,
                       Inches(3.39), Inches(bar_y + 0.12), Inches(9.11), Inches(0.30),
                       size=Pt(10), color=WHITE)

    def save(self, filepath: str | Path) -> Path:
        path = Path(filepath)
        self.prs.save(str(path))
        return path


# ---------------------------------------------------------------------------
# Generator (public API)
# ---------------------------------------------------------------------------

class PptxReleaseNotesGenerator:
    """Generates SE-branded PPTX release notes from a JSON spec file."""

    def __init__(self, settings: Settings) -> None:
        self._settings = settings

    def generate(self, spec_path: str | Path) -> str:
        """
        Build PPTX from a JSON spec file.
        Returns the output file path.
        """
        spec_path = Path(spec_path)
        if not spec_path.exists():
            raise FileNotFoundError(f"Spec file not found: {spec_path}")

        raw = json.loads(spec_path.read_text(encoding="utf-8"))
        content = self._parse_content(raw)
        return self._build_and_write(content)

    def generate_from_content(self, content: PptxContent) -> str:
        """Build PPTX from a PptxContent dataclass directly."""
        return self._build_and_write(content)

    # -- internals ----------------------------------------------------------

    def _parse_content(self, raw: dict) -> PptxContent:
        """Convert raw JSON dict into PptxContent dataclass."""
        exec_summary = [
            ExecutiveSummaryItem(
                title=item.get("title", ""),
                description=item.get("description", ""),
            )
            for item in raw.get("executive_summary", [])
        ]

        features = [
            FeatureSlide(
                title=f.get("title", ""),
                subtitle=f.get("subtitle", ""),
                bullets=[
                    BulletItem(
                        text=b.get("text", "") if isinstance(b, dict) else b,
                        detail=b.get("detail", "") if isinstance(b, dict) else "",
                    )
                    for b in f.get("bullets", [])
                ],
                screenshot_caption=f.get("screenshot_caption", ""),
            )
            for f in raw.get("features", [])
        ]

        df_raw = raw.get("defect_fixes", {})
        defect_fixes = DefectFixesData(
            subtitle=df_raw.get("subtitle", ""),
            fixes=df_raw.get("fixes", []),
        )

        key_benefits = [
            KeyBenefit(
                number=kb.get("number", ""),
                category=kb.get("category", ""),
                title=kb.get("title", ""),
                description=kb.get("description", ""),
            )
            for kb in raw.get("key_benefits", [])
        ]

        return PptxContent(
            product_name=raw.get("product_name", ""),
            version_display=raw.get("version_display", ""),
            version_code=raw.get("version_code", ""),
            release_date=raw.get("release_date", ""),
            executive_summary=exec_summary,
            features=features,
            defect_fixes=defect_fixes,
            key_benefits=key_benefits,
        )

    def _build_and_write(self, content: PptxContent) -> str:
        """Build the full presentation and write to disk."""
        builder = _SlideBuilder()

        # Slide 1: Title
        builder.build_title_slide(content)

        # Slide 2: Executive Summary
        if content.executive_summary:
            builder.build_executive_summary(content)

        # Slides 3-N: Feature details
        for feature in content.features:
            builder.build_feature_slide(feature, content.version_code)

        # Slide N+1: Defect Fixes
        if content.defect_fixes.fixes:
            builder.build_defect_fixes_slide(content.defect_fixes, content.version_code)

        # Slide N+2: Key Benefits
        if content.key_benefits:
            builder.build_key_benefits_slide(content)

        # Write output
        safe_version = sanitize_filename(content.version_code)
        date_str = today_str(self._settings.output.date_format)
        filename = f"{safe_version}_release_notes_{date_str}.pptx"
        out_dir = ensure_dir(self._settings.paths.output_full_release_notes)
        out_path = builder.save(out_dir / filename)
        return str(out_path)
